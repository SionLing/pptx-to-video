import os
import collections
import collections.abc
from pptx import Presentation

# Get pptx notes
def get_pptx_notes(pptx_file):
    prs = Presentation(pptx_file)
    notes = []
    for i,slide in enumerate(prs.slides):
        text = slide.notes_slide.notes_text_frame.text
        notes.append({'text':text})
    return notes

# use convert command line to convert pptx slides to JPG files
def convert_pptx_to_jpg(pptx_file, output_dir):
    cmdline = f'convert -density 144 -quality 100 {pptx_file} {output_dir}/%d.jpg'
    # print(cmdline)
    os.system(cmdline)
    pass

# Export pptx to images and texts
def export_pptx(pptx_file, output_dir):
    convert_pptx_to_jpg(pptx_file, output_dir)
    notes = get_pptx_notes(pptx_file)
    # combine notes and images
    for i,note in enumerate(notes):
        note['image'] = f'{output_dir}/{i}.jpg'
    return notes
